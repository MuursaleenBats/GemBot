import os
import google.generativeai as genai
from AppOpener import open as app_open
from fuzzywuzzy import process
import winapps
import sys
from flask import Flask, request, jsonify
import queue
from flask_cors import CORS
import speech_recognition as sr
import pyttsx3
import json
import subprocess
from pywinauto import Application, timings, mouse
from pywinauto import Desktop, findwindows
from pywinauto.keyboard import send_keys
import winreg
import os
from sentence_transformers import SentenceTransformer
import re
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import win32gui
import win32process
import win32con
import win32api
import psutil
import time
import logging
from installer.app_install import AppInstaller
from AutoWin.automator import automater
import threading
import tkinter as tk
from tkinter import filedialog, simpledialog
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt


logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
tk_queue = queue.Queue()
genai.configure(api_key=os.environ["API_KEY"])
g_model = genai.GenerativeModel('gemini-1.5-pro')
c_model = genai.GenerativeModel('gemini-1.5-pro',  generation_config={"response_mime_type": "application/json"})
status_queue = queue.Queue()
app = Flask(__name__)
CORS(app)
status_lock = threading.Lock()
latest_status = "No updates"
# Event to signal when new status is available
status_event = threading.Event()
installer = AppInstaller()

model = SentenceTransformer('all-MiniLM-L6-v2')  # Load the model once

global is_blind_mode
is_blind_mode = False

def process_tk_queue():
    while True:
        try:
            func, args, kwargs = tk_queue.get(timeout=1)
            func(*args, **kwargs)
            tk_queue.task_done()
        except queue.Empty:
            break

def get_app_path(app_name):
    try:
        key = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths")
        for i in range(winreg.QueryInfoKey(key)[0]):
            try:
                subkey_name = winreg.EnumKey(key, i)
                if subkey_name.lower() == app_name.lower() or subkey_name.lower() == app_name.lower() + '.exe':
                    subkey = winreg.OpenKey(key, subkey_name)
                    path, _ = winreg.QueryValueEx(subkey, "")
                    return os.path.normpath(path)
            except WindowsError:
                continue

        if app_name.lower() == "explorer" or app_name.lower() == "explorer.exe":
            key_2 = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Winlogon")
            explorer_path, _ = winreg.QueryValueEx(key_2, "Shell")
            return os.path.normpath(explorer_path)

    except WindowsError as e:
        print(f"Failed to get {app_name} path. Error: {e}")
    return None

is_processing_command = False
last_command_time = 0
debounce_interval = 5  # Time interval in seconds to debounce commands

def listen_for_keyword():
    global latest_status, is_processing_command, last_command_time, is_blind_mode
    recognizer = sr.Recognizer()
    while True:
        with sr.Microphone() as source:
            print("Listening for 'Gemini'...")
            audio = recognizer.listen(source)
        try:
            text = recognizer.recognize_google(audio).lower()
            print(f"Heard: {text}")
            if is_blind_mode:
                speak_text(f"Heard: {text}")
            if "gemini" in text and not is_processing_command:
                current_time = time.time()
                if current_time - last_command_time >= debounce_interval:
                    #message = "Keyword 'Gemini' detected!"
                    # if is_blind_mode:
                    #     speak_text(message)
                    # print(message)
                    is_processing_command = True
                    with status_lock:
                        latest_status = "Listening for command..."
                    status_event.set()
                    
                    command = listen_for_command(timeout=8)
                    
                    if command:
                        result = process_command(command)
                        print(f"Command result: {result}")
                        with status_lock:
                            latest_status = json.dumps({
                                "command": command,
                                "response": result
                            })
                        status_event.set()
                    else:
                        message = "No command received within 5 seconds."
                        if is_blind_mode:
                            speak_text(message)
                        print(message)
                        with status_lock:
                            latest_status = "No command received"
                        status_event.set()
                    
                    last_command_time = time.time()
                    is_processing_command = False
                else:
                    message = "Command debounced. Waiting for debounce interval to pass."
                    if is_blind_mode:
                        speak_text(message)
                    print(message)
        except sr.UnknownValueError:
            with status_lock:
                latest_status = "Could not understand audio"
            status_event.set()
        except sr.RequestError as e:
            message = f"Could not request results from Google Speech Recognition service; {e}"
            if is_blind_mode:
                speak_text(message)
            print(message)
            with status_lock:
                latest_status = f"Error: {e}"
            status_event.set()

def speak_text(text):
    engine = pyttsx3.init()
    
    # Get available voices
    voices = engine.getProperty('voices')
    
    # Set voice to female (usually the second voice in the list)
    engine.setProperty('voice', voices[1].id)
    
    # Adjust rate (speed) for more natural sound (default is 200)
    engine.setProperty('rate', 150)
    
    # Adjust volume (0.0 to 1.0)
    engine.setProperty('volume', 0.8)
    
    engine.say(text)
    engine.runAndWait()

def listen_for_command(timeout=8):
    global is_blind_mode
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        if is_blind_mode:
            speak_text(f"Listening for command")
        print(f"Listening for command for {timeout} seconds...")
        try:
            audio = recognizer.listen(source, timeout=timeout)
            command = recognizer.recognize_google(audio)
            print(f"Recognized command: {command}")
            if is_blind_mode:
                speak_text(f"Recognized command: {command}")
            return command
        except sr.WaitTimeoutError:
            message = "Listening timed out"
            if is_blind_mode:
                speak_text(message)
            print(message)
            return None
        except sr.UnknownValueError:
            message = "Could not understand audio"
            if is_blind_mode:
                speak_text(message)
            print(message)
            return None
        except sr.RequestError as e:
            message = f"Could not request results from Google Speech Recognition service; {e}"
            if is_blind_mode:
                speak_text(message)
            print(message)
            return None
        
def start_application(app_name):
    def get_installed_apps():
        apps = list(winapps.list_installed())
        app_names = [app.name for app in apps]
        return app_names

    def suggest_apps(app_name, app_names, threshold=70):
        suggestions = process.extract(app_name, app_names, limit=5)
        return [suggestion for suggestion, score in suggestions if score >= threshold]

    def open_app(app_names):
        installed_apps = get_installed_apps()

        for app_name in app_names:
            try:
                app_open(app_name, match_closest=True)
                return f"I've started {app_name} for you."
            except Exception as e:
                print(f"appopener failed to open {app_name}: {e}")
                suggestions = suggest_apps(app_name, installed_apps)

                if suggestions:
                    print(f"Application '{app_name}' not found. Did you mean one of these?")
                    for suggestion in suggestions:
                        print(f" - {suggestion}")
                    suggested_name = suggestions[0]
                    user_confirmation = input(f"Do you want to open '{suggested_name}' instead? (yes/no): ").strip().lower()
                    if user_confirmation == 'yes':
                        open_app([suggested_name])
                else:
                    print("Attempting to open the app using winapps...")
                    app = list(winapps.search_installed(app_name))

                    if app:
                        app_info = app[0]
                        uninstall_string = app_info.uninstall_string

                        try:
                            subprocess.Popen(uninstall_string)
                            print(f"Opening {app_info.name} using uninstall string...")
                        except Exception as e:
                            print(f"Error opening {app_info.name} using uninstall string: {e}")
                            return f"I'm sorry, I couldn't find {app_name} on your system. Could you check the spelling or try another app?"
                    else:
                        return f"I'm sorry, I couldn't find {app_name} on your system. Could you check the spelling or try another app?"
                        
    return open_app([app_name])

def install_application(code):
    try:
        result = subprocess.run(f'powershell -Command "{code}"', shell=True, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return "The installation process has started. It may take a few minutes to complete."
    except subprocess.CalledProcessError as e:
        return f"I encountered an issue while trying to install the application. Could you check your internet connection and try again?"

def open_website_in_chrome(url):
    chrome_path = get_app_path("chrome")
    if chrome_path:
        try:
            app = Application().start(f'"{chrome_path}" {url}')
            return f"I've opened {url} in Google Chrome for you."
        except Exception as e:
            return f"I had trouble opening the website. Is Chrome running properly?"
    else:
        return "I couldn't find Google Chrome on your system. Is it installed?"

def send_to_gemini(command, model):
    try:
        response = model.generate_content(command)
        if hasattr(response, 'text'):
            return response.text
        elif isinstance(response, str):
            return response
        else:
            print(f"Unexpected response type from Gemini: {type(response)}")
            return str(response)
    except Exception as e:
        print(f"Error sending command to Gemini: {e}")
        return None

def list_windows():
    def winEnumHandler(hwnd, ctx):
        if win32gui.IsWindowVisible(hwnd):
            title = win32gui.GetWindowText(hwnd)
            if title:
                try:
                    _, pid = win32process.GetWindowThreadProcessId(hwnd)
                    process = psutil.Process(pid)
                    ctx.append(f"{title} (Process: {process.name()})")
                except:
                    ctx.append(title)
    windows = []
    win32gui.EnumWindows(winEnumHandler, windows)
    return windows

def find_window(partial_name):
    partial_name = partial_name.lower()
    matching_windows = []
    for window in list_windows():
        if partial_name in window.lower():
            matching_windows.append(window)
    return matching_windows

def close_window_function(partial_name):
    matching_windows = find_window(partial_name)
    if matching_windows:
        if len(matching_windows) == 1:
            full_name = matching_windows[0]
            title = full_name.split(" (Process: ")[0]
            process_name = full_name.split("(Process: ")[1][:-1]  # Remove the last ')'

            def enum_windows_callback(hwnd, result):
                if win32gui.IsWindowVisible(hwnd) and title.lower() in win32gui.GetWindowText(hwnd).lower():
                    _, pid = win32process.GetWindowThreadProcessId(hwnd)
                    if psutil.Process(pid).name().lower() == process_name.lower():
                        result.append(hwnd)
                return True

            result = []
            win32gui.EnumWindows(enum_windows_callback, result)

            if result:
                hwnd = result[0]
                try:
                    win32gui.PostMessage(hwnd, win32con.WM_CLOSE, 0, 0)
                    return f"I've closed the window '{full_name}' for you."
                except Exception as e:
                    return f"I had trouble closing the window. Could you try closing it manually?"
            else:
                return f"I found multiple windows matching '{partial_name}'. Could you be more specific?"
        else:
            return f"I couldn't find any windows matching '{partial_name}'. Is the window open?"

def generate_and_save_code(user_input, model):
    language = extract_language(user_input)

    prompt = f"Generate {language} code for: {user_input}. Provide only the code, without explanations."
    generated_code = send_to_gemini(prompt, model)

    if generated_code:
        extension = get_file_extension(language)

        # Remove any potential markdown code block syntax
        generated_code = re.sub(r'^```.*?\n|```$', '', generated_code, flags=re.DOTALL).strip()

        # Create a root window
        root = tk.Tk()

        # Ask for file name
        file_name = "generated_code"

        if not file_name.endswith(extension):
            file_name += extension

        # Ask for save location
        file_path = filedialog.asksaveasfilename(
            defaultextension=extension,
            filetypes=[(f"{language.capitalize()} files", f"*{extension}"), ("All files", "*.*")],
            title="Save Generated Code",
            initialfile=file_name
        )

        if not file_path:  # User cancelled
            print("File save cancelled.")
            return None

        try:
            with open(file_path, 'w') as file:
                file.write(generated_code)
            print(f"Generated code saved to {file_path}")
        except IOError as e:
            print(f"Error saving file: {e}")
            return None

        # Execute the code if it's Python
        if language.lower() == 'python':
            print("\nExecuting the generated Python code:")
            try:
                # Use subprocess to run the Python script
                result = subprocess.run([sys.executable, file_path], capture_output=True, text=True, timeout=10)
                print("Output:")
                print(result.stdout)
                result_output = f"Code: {generated_code} \nOutput: {result.stdout}"
                if result.stderr:
                    print("Errors:")
                    print(result.stderr)
                    result_output += f"\nErrors: {result.stderr}"
            except subprocess.TimeoutExpired:
                print("Execution timed out after 10 seconds.")
                result_output = "Execution timed out after 10 seconds."
            except Exception as e:
                print(f"An error occurred while executing the code: {e}")
                result_output = f"An error occurred while executing the code: {e}"
        else:
            print(f"\nNote: Automatic execution is only supported for Python. The {language} code has been saved but not executed.")
            result_output = f"Generated {language} code has been saved."

        return result_output
    else:
        print("Failed to generate code.")
        return None

def extract_language(text):
    languages = ['python', 'java', 'javascript', 'c++', 'c#', 'ruby']
    for lang in languages:
        if lang in text.lower():
            return lang
    return 'python'  # Default to Python if no language is specified

def get_file_extension(language):
    extensions = {
        'python': '.py',
        'java': '.java',
        'javascript': '.js',
        'c++': '.cpp',
        'c#': '.cs',
        'ruby': '.rb'
    }
    return extensions.get(language, '.txt')

def process_command(command, max_retries=1):
    global is_blind_mode
    if not command.strip():
        return "I didn't catch that. Could you please repeat your command?"
    if "i am blind" in command.lower():
        is_blind_mode = True
        response = "Blind mode activated. All future prompts and responses will be read aloud."
        speak_text(response)
        return response

    if is_blind_mode:
        speak_text(f"You said: {command}")

    prompt = f"""
    Given the user command: "{command}"
    Generate a JSON response with the appropriate UI action and parameters.
    Possible actions are:
    1. close_window: Requires a "window_name" parameter
    2. list_windows: No parameters required
    3. interact_with_control: Requires "window_name", "control_type", "control_name", and "action" parameters
    4. navigate_in_browser: Requires "browser_name" and "url" parameters
    5. file_explorer_operation: Requires "operation" and additional parameters based on the operation
    6. list_processes: No parameters required
    7. get_process_info: Requires a "pid" parameter
    8. open_website_in_chrome: Requires a "url" parameter
    9. start_application: Requires an "app_name" parameter
    10. install_application: Requires an "app_name" parameter
    11. generate_and_save_code: Requires a "language" and "code_description" parameter
    12. generate_and_save_word_document: Requires a "content" parameter
    13. generate_powerpoint: Requires "title"
    14. no_action: Use this if no UI action is needed
    15. process_request: Use this if none of the above actions are suitable. Requires a "request" parameter with the original user command.

    Example response formats:
    {{"action": "close_window", "params": {{"window_name": "Firefox"}}}}
    {{"action": "list_windows"}}
    {{"action": "interact_with_control", "params": {{"window_name": "Notepad", "control_type": "Edit", "control_name": "Text Editor", "action": "type_keys", "value": "Hello, World!"}}}}
    {{"action": "navigate_in_browser", "params": {{"browser_name": "Chrome", "url": "https://www.example.com"}}}}
    {{"action": "file_explorer_operation", "params": {{"operation": "rename_files", "folder_path": "C:/Users/YourName/Documents", "file_pattern": "*.txt", "new_name": "renamed_{{index}}.txt"}}}}
    {{"action": "list_processes"}}
    {{"action": "get_process_info", "params": {{"pid": 1234}}}}
    {{"action": "open_website_in_chrome", "params": {{"url": "https://www.example.com"}}}}
    {{"action": "start_application", "params": {{"app_name": "winword.exe"}}}}
    {{"action": "install_application", "params": {{"app_name": "Chrome"}}}}
    {{"action": "generate_and_save_code", "params": {{"language": "python", "code_description": "A function to calculate fibonacci numbers"}}}}
    {{"action": "generate_powerpoint", "params": {{"title": "Automated Presentation"}}}}
    {{"action": "no_action", "response": "Hello! How can I assist you with UI automation today?"}}
    {{"action": "process_request", "params": {{"request": "Original user command that doesn't fit other actions"}}}}

    For the generate_and_save_word_document action, provide a detailed, well-structured content for a Word document. The content should be rich, informative, and well-organized. Include appropriate headings, subheadings, and paragraphs. The content should be at least 500 words long and cover the topic comprehensively.

    Example response format for generate_and_save_word_document:
    {{"action": "generate_and_save_word_document", "params": {{"content": "# Title of the Document\\n\\n## Introduction\\n[Detailed introduction paragraph]\\n\\n## Main Section 1\\n[Comprehensive content for section 1]\\n\\n### Subsection 1.1\\n[Detailed information for subsection 1.1]\\n\\n### Subsection 1.2\\n[Detailed information for subsection 1.2]\\n\\n## Main Section 2\\n[Comprehensive content for section 2]\\n\\n### Subsection 2.1\\n[Detailed information for subsection 2.1]\\n\\n### Subsection 2.2\\n[Detailed information for subsection 2.2]\\n\\n## Conclusion\\n[Detailed concluding paragraph]\\n\\n"}}}}

    Ensure that the generated content is relevant to the user's command, well-structured, and provides valuable information on the topic. The response should be a valid JSON object.
    If none of the specific actions are suitable, use the "process_request" action with the original user command as the request parameter.
    """

    retries = 0
    while retries < max_retries:
        try:
            if "code" in prompt:
                response = g_model.generate_content(prompt)
                logging.info(f"Raw response from Gemini: {response.text}")

                json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
                if not json_match:
                    return "I'm having trouble understanding that request. Could you rephrase it?"

                json_str = json_match.group(0)

            response = c_model.generate_content(prompt)
            logging.info(f"Raw response from Gemini: {response.text}")

            json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
            if not json_match:
                return "I'm having trouble processing that request. Could you try again?"

            json_str = json_match.group(0)

            json_str = json_str.replace('\n', '\\n').replace('\r', '\\r').replace('\t', '\\t')

            try:
                action_data = json.loads(json_str)
            except json.JSONDecodeError as json_err:
                logging.error(f"JSON parsing error: {str(json_err)}")
                logging.error(f"Problematic JSON string: {json_str}")
                return "I encountered an issue while processing your request. Could you try rephrasing it?"

            if not isinstance(action_data, dict):
                return "I'm having trouble interpreting your request. Could you provide more details?"

            if "action" not in action_data:
                return "I'm not sure what action to take. Could you be more specific?"

            if action_data["action"] == "no_action":
                return action_data.get("response", "I'm not sure how to help with that. Could you clarify?")

            result = execute_ui_action(action_data["action"], action_data.get("params", {}))

            if is_blind_mode:
                speak_text(f"Response: {result}")

            return result
        except Exception as e:
            logging.error(f"Error in process_command: {str(e)}")
            return "I encountered an unexpected issue. Could you try your request again?"
        
def generate_powerpoint(title):
    prompt = f"""
    Given the user command: "Generate a PowerPoint presentation with the title '{title}'"
    Generate a JSON response with the appropriate parameters for creating the PowerPoint presentation.
    The JSON response should include the following fields:
    - "title": The title of the presentation
    - "subtitle": The subtitle of the presentation
    - "agenda_items": An array of at least 5 agenda items for the presentation, each with a 2-line description
    - "filename": The filename for the PowerPoint file

    Example response format:
    {{
        "title": "Automated Presentation",
        "subtitle": "Created using python-pptx",
        "agenda_items": [
            {{
                "item": "Introduction",
                "description": "This slide introduces the topic. It provides an overview of what will be covered."
            }},
            {{
                "item": "Main Topic 1",
                "description": "This slide covers the first main topic. It discusses key points and relevant information."
            }},
            {{
                "item": "Main Topic 2",
                "description": "This slide covers the second main topic. It provides additional details and insights."
            }},
            {{
                "item": "Main Topic 3",
                "description": "This slide covers the third main topic. It includes important facts and data."
            }},
            {{
                "item": "Conclusion",
                "description": "This slide summarizes the main points. It provides a final overview and closing remarks."
            }}
        ],
        "filename": "automated_presentation.pptx"
    }}
    """

    try:
        response = g_model.generate_content(prompt)
        print(response)
        json_match = re.search(r'```json\s*(.*?)\s*```', response.text, re.DOTALL)
        if json_match:
            json_str = json_match.group(1)
        else:
            json_str = response.text

        action_data = json.loads(json_str)

        if not isinstance(action_data, dict):
            raise ValueError(f"Invalid response format. Expected a dictionary, got: {type(action_data)}")

        title = action_data.get("title", "Automated Presentation")
        subtitle = action_data.get("subtitle", "Created using python-pptx")
        agenda_items = action_data.get("agenda_items", [
            {"item": "Introduction", "description": "This slide introduces the topic. It provides an overview of what will be covered."},
            {"item": "Main Topic 1", "description": "This slide covers the first main topic. It discusses key points and relevant information."},
            {"item": "Main Topic 2", "description": "This slide covers the second main topic. It provides additional details and insights."},
            {"item": "Main Topic 3", "description": "This slide covers the third main topic. It includes important facts and data."},
            {"item": "Conclusion", "description": "This slide summarizes the main points. It provides a final overview and closing remarks."}
        ])
        filename = action_data.get("filename", "automated_presentation.pptx")

        prs = Presentation()

        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle

        # Add a content slide for the agenda
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Agenda"

        tf = body_shape.text_frame
        tf.text = ""  # Clear any existing text

        for item in agenda_items:
            p = tf.add_paragraph()
            p.text = item["item"]
            p.level = 1

        # Add a separate slide for each agenda item with description
        for item in agenda_items:
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = item["item"]

            tf = body_shape.text_frame
            tf.text = item["description"]

        # Save the presentation
        # prs.save(filename)

        # Save the presentation
        def save_presentation():
            root = tk.Tk()

            file_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint presentations", "*.pptx"), ("All files", "*.*")],
                title="Save Generated PowerPoint Presentation",
                initialfile=filename
            )

            if not file_path:  # User cancelled
                return "You've cancelled saving the presentation. Let me know if you want to save it later."
                return None

            prs.save(file_path)
            print(f"Generated PowerPoint presentation saved to {file_path}")
            return f"I've saved the PowerPoint presentation for you. You can find it at {file_path}"

        # Add the save_presentation function to the queue
        tk_queue.put((save_presentation, (), {}))
        # Process the queue in the main thread
        process_tk_queue()

        return "PowerPoint presentation generation initiated."
    except Exception as e:
        logging.error(f"Error in generate_powerpoint: {str(e)}")
        return "I had trouble creating the PowerPoint presentation. Is Microsoft PowerPoint installed on your system?"


def generate_and_save_word_document(content):
    try:
        # Create a new Document
        doc = Document()

        # Parse the markdown-like content and add it to the document
        lines = content.split('\n')
        for line in lines:
            if line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('### '):
                doc.add_heading(line[4:], level=3)
            elif line.strip():
                doc.add_paragraph(line)

        # Save the document
        def save_document():
            root = tk.Tk()

            # Ask for file name
            file_name = "generated_doc"

            if not file_name.endswith(".docx"):
                file_name += ".docx"

            # Ask for save location
            file_path = filedialog.asksaveasfilename(
                defaultextension=".docx",
                filetypes=[("Word documents", "*.docx"), ("All files", "*.*")],
                title="Save Generated Word Document",
                initialfile=file_name
            )

            if not file_path:  # User cancelled
                print("File save cancelled.")
                return None

            doc.save(file_path)
            print(f"Generated Word document saved to {file_path}")
            return f"Generated Word document saved to {file_path}"

        # Add the save_document function to the queue
        tk_queue.put((save_document, (), {}))
        # Process the queue in the main thread
        process_tk_queue()

        return "Word document generation initiated."
    except Exception as e:
        logging.error(f"Error generating and saving Word document: {str(e)}")
        return "I encountered an issue while creating the Word document. Do you have Microsoft Word installed and enough disk space?"


def interact_with_control(params):
    window_name = params.get("window_name")
    control_type = params.get("control_type")
    control_name = params.get("control_name")
    action = params.get("action")
    value = params.get("value")

    if not all([window_name, control_type, control_name, action]):
        return "I need more information about the control you want to interact with. Could you provide more details?"

    try:
        app = Application(backend="uia").connect(title=window_name, timeout=10)
        window = app.window(title=window_name)
        control = getattr(window, control_type)(title=control_name)

        if action == "type_keys":
            control.type_keys(value, with_spaces=True)
        elif action == "click":
            control.click_input()
        elif action == "double_click":
            control.double_click_input()
        elif action == "right_click":
            control.right_click_input()
        elif action == "scroll":
            control.scroll(value)
        else:
            return f"Unknown action for control: {action}"

        return f"Successfully performed {action} on {control_type} '{control_name}' in {window_name}"
    except Exception as e:
        logging.error(f"Error in interact_with_control: {str(e)}")
        return f"I had trouble interacting with the {control_type} in {window_name}. Could you check if the window is open and try again?"

def navigate_in_browser(browser_name, url, search_query=None):
    if not browser_name or not url:
        return "I need both a browser name and a URL to navigate. Could you provide both?"

    try:
        print("Listing all windows for debugging:")
        windows = list_windows()
        for window in windows:
            print(f"Window: {window}")

        app = Application(backend="uia").connect(title_re=f".*{browser_name}.*", found_index=0, timeout=10)
        all_windows = app.windows()

        # Try to find the main browser window
        browser_window = None
        for win in all_windows:
            win_title = win.window_text().lower()
            logging.info(f"Checking window: {win.window_text()}")
            if browser_name.lower() in win_title:
                browser_window = win
                break

        if not browser_window:
            raise Exception("No suitable browser window found")

        logging.info(f"Connected to {browser_name} window: {browser_window.window_text()}")

        # Open a new tab first
        browser_window.set_focus()
        send_keys('^t')  # Ctrl+T to open a new tab
        time.sleep(1)  # Wait for the new tab to open

        # Navigate to Google if the search query is provided
        if search_query:
            send_keys('^l')  # Ctrl+L to focus on the address bar
            time.sleep(0.5)
            send_keys('^a')  # Select all existing text
            send_keys('{BACKSPACE}')  # Clear the address bar
            send_keys('https://www.google.com{ENTER}')
            time.sleep(2)  # Wait for Google to load

            # Simulate typing the search query and pressing Enter
            send_keys(search_query + '{ENTER}')
        else:
            # If no search query is provided, just navigate to the URL
            send_keys('^l')  # Ctrl+L to focus on the address bar
            time.sleep(0.5)
            send_keys('^a')  # Select all existing text
            send_keys('{BACKSPACE}')  # Clear the address bar
            send_keys(url + '{ENTER}')
            time.sleep(2)  # Wait for the page to load

        return f"Opened a new tab and navigated to {url} in {browser_name}"
    except Exception as e:
        logging.error(f"Error in navigate_in_browser: {str(e)}")
        return f"I couldn't navigate to the website in {browser_name}. Is the browser open and connected to the internet?"

def get_current_file_explorer_path():
    try:
        hwnds = find_window("File Explorer")
        if not hwnds:
            return None

        for hwnd in hwnds:
            app = Application(backend="uia").connect(handle=hwnd)
            explorer_window = app.top_window()
            address_bar = explorer_window.child_window(auto_id="Address Band Root", control_type="ToolBar")
            address_edit = address_bar.child_window(control_type="Edit")
            path = address_edit.get_value()
            if path:
                return path

        return None
    except Exception as e:
        logging.error(f"Error getting current File Explorer path: {str(e)}")
        return None

def file_explorer_operation(params):
    operation = params.get("operation")
    folder_path = params.get("folder_path")

    if not operation:
        return "I'm not sure what operation you want to perform in File Explorer. Could you specify?"

    try:
        if not folder_path:
            folder_path = get_current_file_explorer_path()
            if not folder_path:
                folder_path = input("Folder path not found. Please enter the folder path: ")
            if not os.path.exists(folder_path):
                return f"Folder not found: {folder_path}"

        os.startfile(folder_path)
        time.sleep(1)

        app = Application(backend="uia").connect(title_re=".*File Explorer.*", timeout=10)
        explorer_window = app.top_window()

        if operation == "rename_files":
            return rename_files(explorer_window, params)
        elif operation == "select_files":
            return select_files(explorer_window, params)
        elif operation == "navigate_to_folder":
            return navigate_to_folder(explorer_window, params)
        else:
            return f"Unknown file explorer operation: {operation}"
    except Exception as e:
        logging.error(f"Error in file_explorer_operation: {str(e)}")
        return f"I encountered an issue while trying to perform the {operation} in File Explorer. Could you check if the folder exists and try again?"

def rename_files(explorer_window, params):
    file_pattern = params.get("file_pattern")
    new_name = params.get("new_name")

    if not file_pattern or not new_name:
        return "I need both a file pattern and a new name to rename files. Could you provide both?"

    try:
        folder_path = get_current_file_explorer_path()
        if not folder_path:
            return "Could not retrieve the current File Explorer path"

        files = [f for f in os.listdir(folder_path) if re.match(file_pattern, f)]
        if not files:
            return f"No files matching the pattern '{file_pattern}' found in '{folder_path}'"

        for index, file in enumerate(files, start=1):
            new_file_name = new_name.replace("{index}", str(index))
            os.rename(os.path.join(folder_path, file), os.path.join(folder_path, new_file_name))

        return f"Renamed files matching '{file_pattern}' to '{new_name}' in '{folder_path}'"
    except Exception as e:
        logging.error(f"Error in rename_files: {str(e)}")
        return "I had trouble renaming the files. Are you sure the files exist and you have permission to rename them?"

def select_files(explorer_window, params):
    file_pattern = params.get("file_pattern")

    if not file_pattern:
        return "Missing file pattern"

    try:
        explorer_window.set_focus()
        send_keys(f'{file_pattern}')
        time.sleep(0.5)
        send_keys('^a')

        return f"Selected files matching '{file_pattern}'"
    except Exception as e:
        logging.error(f"Error in select_files: {str(e)}")
        return f"Error selecting files"

def navigate_to_folder(explorer_window, params):
    target_folder = params.get("target_folder")

    if not target_folder:
        return "I need to know which folder you want to navigate to. Could you specify the folder path?"

    try:
        explorer_window.set_focus()
        send_keys('^l')  # Ctrl+L to focus on the address bar
        time.sleep(0.5)
        send_keys('^a')  # Select all existing text
        send_keys('{BACKSPACE}')  # Clear the address bar
        send_keys(target_folder + '{ENTER}')
        time.sleep(2)  # Wait for the folder to load

        return f"I've navigated to the folder: {target_folder}"
    except Exception as e:
        logging.error(f"Error in navigate_to_folder: {str(e)}")
        return f"I had trouble navigating to {target_folder}. Does this folder exist on your system?"

def list_processes():
    processes = []
    for proc in psutil.process_iter(['pid', 'name', 'status']):
        try:
            processes.append(f"Process: {proc.info['name']}, Status: {proc.info['status']}")
        except (psutil.NoSuchProcess, psutil.AccessDenied, psutil.ZombieProcess):
            pass
    return "Here are the currently running processes:\n" + "\n".join(processes[:20]) + "\n(Showing first 20 processes)"

def get_process_info(pid):
    try:
        handle = win32api.OpenProcess(win32con.PROCESS_QUERY_INFORMATION | win32con.PROCESS_VM_READ, False, pid)
        exe = win32process.GetModuleFileNameEx(handle, 0)
        return exe
    except:
        return None

def get_detailed_process_info(pid):
    try:
        pid = int(pid)
        process = psutil.Process(pid)
        exe_path = get_process_info(pid)
        return f"PID: {pid}, Name: {process.name()}, Exe: {exe_path}, Status: {process.status()}"
    except Exception as e:
        logging.error(f"Error in get_detailed_process_info: {str(e)}")
        return f"I couldn't retrieve information for the process with PID {pid}. Is the process still running?"
    
def execute_ui_action(action, params):
    try:
        if action == "close_window":
            return close_window_function(params.get("window_name"))
        elif action == "interact_with_control":
            return interact_with_control(params)
        elif action == "navigate_in_browser":
            return navigate_in_browser(params.get("browser_name"), params.get("url"), params.get("search_query"))
        elif action == "file_explorer_operation":
            return file_explorer_operation(params)
        elif action == "list_processes":
            return "\n".join(list_processes())
        elif action == "get_process_info":
            return get_detailed_process_info(params.get("pid"))
        elif action == "open_website_in_chrome":
            return open_website_in_chrome(params.get("url"))
        elif action == "start_application":
            return start_application(params.get("app_name"))
        elif action == "install_application":
            return installer.install_app(params.get("app_name"))
        elif action == "generate_and_save_code":
            return generate_and_save_code(params.get("language") + params.get("code_description"), g_model)
        elif action == "generate_and_save_word_document":
             return generate_and_save_word_document(params.get("content", ""))
        elif action == "generate_powerpoint":
            return generate_powerpoint(params.get("title"))
        elif action == "process_request":
            return automater.process_request(params.get("request"))
        else:
            return f"Unknown action: {action}"
    except Exception as e:
        logging.error(f"Error in execute_ui_action: {str(e)}")
        return "I had trouble performing that action. Could you try again or rephrase your request?"
    
def extract_url(text):
    url_pattern = re.compile(r"[a-zA-Z0-9.-]+\.(com|org|net|edu|gov)")
    match = url_pattern.search(text)
    if match:
        return match.group(0)
    return None

def extract_app_name(text):
    start_match = re.search(r'\bstart\b\s+(\w+)', text)
    install_match = re.search(r'\binstall\b\s+(\w+)', text)
    close_match = re.search(r'\bclose\b\s+(\w+)', text)
    if start_match:
        return start_match.group(1)
    elif install_match:
        return install_match.group(1)
    elif close_match:
        return close_match.group(1)
    return None

@app.route('/command', methods=['POST'])
def receive_command():
    data = request.json
    if not data or 'command' not in data:
        return jsonify({"error": "No command provided"}), 400

    command = data['command']
    result = process_command(command)
    return jsonify({"result": result})

@app.route('/status', methods=['GET'])
def get_status():
    global latest_status
    timeout = 3  # Time to wait for new status (in seconds)
    status_event.wait(timeout)  # Wait for new status or timeout
    status_event.clear()  # Reset the event
    print(latest_status)
    try:
        status_data = json.loads(latest_status)
        latest_status = "No updates"
        return jsonify(status_data)
    except:
        return jsonify({"status": latest_status})
    
def create_hidden_root():
        root = tk.Tk()
        root.withdraw() # This hides the window
        root.attributes("-topmost", True)  
        root.mainloop()

def run_flask():
    app.run(port=5000, debug=True, use_reloader=False, threaded=True)

if __name__ == "__main__":
    # Start the voice recognition thread
    voice_thread = threading.Thread(target=listen_for_keyword, daemon=True)
    voice_thread.start()

    # Start the Flask server in a separate thread
    flask_thread = threading.Thread(target=run_flask, daemon=True)
    flask_thread.start()

    
    # Keep the main thread alive
    tk_thread = threading.Thread(target=create_hidden_root, daemon=True)
    tk_thread.start()


    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        print("Shutting down...")
